"""
Generate a 5-slide PowerPoint presentation demonstrating React skills learned.
Uses the RRA/React_Presentation.pptx template for consistent branding.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

TEMPLATE = os.path.join(os.path.dirname(__file__), "RRA", "React_Presentation.pptx")
OUTPUT = os.path.join(os.path.dirname(__file__), "React_Skills_Presentation.pptx")

# Colors from template
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
GOLD = RGBColor(0xD4, 0x8B, 0x0A)
LIGHT_BLUE = RGBColor(0x15, 0x90, 0xD0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xDC, 0xE7, 0xF6)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)


def set_font(run, name="Roboto", size=16, bold=False, color=DARK_TEXT):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=16,
                 bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name="Roboto"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, name=font_name, size=font_size, bold=bold, color=color)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=DARK_TEXT, bold_prefix=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.space_before = Pt(6)
        p.space_after = Pt(4)
        p.level = 0

        # Split on first colon for bold prefix
        if bold_prefix and ":" in item:
            prefix, rest = item.split(":", 1)
            run1 = p.add_run()
            run1.text = "\u2022  " + prefix + ":"
            set_font(run1, size=font_size, bold=True, color=color)
            run2 = p.add_run()
            run2.text = rest
            set_font(run2, size=font_size, bold=False, color=color)
        else:
            run = p.add_run()
            run.text = "\u2022  " + item
            set_font(run, size=font_size, bold=False, color=color)

    return txBox


def add_gold_accent_bar(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()
    return shape


def add_blue_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(1, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BLUE
    shape.line.fill.background()
    return shape


def add_section_header(slide, text, top=Inches(0.6)):
    """Add a section header with gold accent bar."""
    add_gold_accent_bar(slide, Inches(0.45), top, Pt(5), Inches(0.7))
    add_text_box(
        slide, Inches(0.7), top, Inches(8), Inches(0.5),
        text, font_size=28, bold=True, color=WHITE, font_name="Calibri"
    )
    add_blue_accent_line(slide, Inches(0.7), top + Inches(0.65), Inches(8))


def add_slide_number(slide, num):
    add_text_box(
        slide, Inches(0.3), Inches(6.6), Inches(0.5), Inches(0.3),
        str(num), font_size=10, color=RGBColor(0x99, 0x99, 0x99)
    )


def create_content_slide(prs, layout_index=1):
    """Create a slide with the dark blue background matching the template."""
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    # Set background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    # Add bottom bar
    bottom_bar = slide.shapes.add_shape(1, 0, Inches(6.1), Inches(12.5), Inches(1.3))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = LIGHT_BG
    bottom_bar.line.fill.background()
    # Gold line
    add_gold_accent_bar(slide, 0, Inches(5.95), Inches(12.5), Pt(4))
    return slide


def main():
    prs = Presentation(TEMPLATE)

    # Remove existing slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId is None:
            rId_attr = list(prs.slides._sldIdLst[0].attrib.keys())
            for attr in rId_attr:
                if 'id' in attr.lower() and 'r' in attr.lower():
                    rId = prs.slides._sldIdLst[0].get(attr)
                    break
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # =========================================================================
    # SLIDE 1: Title Slide
    # =========================================================================
    slide1 = create_content_slide(prs, 0)
    add_gold_accent_bar(slide1, Inches(0.45), Inches(0.8), Pt(5), Inches(1.2))
    add_text_box(
        slide1, Inches(0.85), Inches(0.8), Inches(9), Inches(0.7),
        "React: Key Skills Acquired", font_size=36, bold=True,
        color=WHITE, font_name="Calibri"
    )
    add_text_box(
        slide1, Inches(0.85), Inches(1.6), Inches(9), Inches(0.5),
        "Frontend Development Training", font_size=24,
        bold=False, color=GOLD, font_name="Calibri"
    )
    add_text_box(
        slide1, Inches(3), Inches(4.3), Inches(6), Inches(0.3),
        "Rwanda Revenue Authority (RRA)", font_size=14,
        color=RGBColor(0xBB, 0xCC, 0xDD), font_name="Calibri"
    )
    add_text_box(
        slide1, Inches(0.85), Inches(3.2), Inches(5), Inches(0.3),
        "NINDENKIMANA Blaise  |  Strategy and Risk Analysis Department",
        font_size=12, color=RGBColor(0xBB, 0xCC, 0xDD), font_name="Roboto"
    )
    add_text_box(
        slide1, Inches(0.85), Inches(3.6), Inches(5), Inches(0.3),
        "23rd April 2026",
        font_size=12, color=LIGHT_BLUE, bold=True, font_name="Roboto"
    )
    add_slide_number(slide1, 1)

    # =========================================================================
    # SLIDE 2: What is React?
    # =========================================================================
    slide2 = create_content_slide(prs, 0)
    add_section_header(slide2, "What is React?")

    # Left column - History & Definition
    add_text_box(
        slide2, Inches(0.85), Inches(1.45), Inches(5), Inches(0.3),
        "History & Overview", font_size=15, bold=True, color=GOLD
    )
    add_bullet_list(
        slide2, Inches(0.85), Inches(1.85), Inches(5.2), Inches(1.5),
        [
            "React is an open-source JavaScript library for building user interfaces, created by Jordan Walke at Facebook (now Meta)",
            "First deployed on Facebook's News Feed in 2011, then on Instagram in 2012; open-sourced in May 2013 at JSConf US",
            "Maintained by Meta and a large community of developers; currently at version 19 with continuous innovation",
        ],
        font_size=11, color=WHITE, bold_prefix=False
    )

    add_text_box(
        slide2, Inches(0.85), Inches(3.3), Inches(5), Inches(0.3),
        "Core Capabilities", font_size=15, bold=True, color=GOLD
    )
    add_bullet_list(
        slide2, Inches(0.85), Inches(3.7), Inches(5.2), Inches(2.0),
        [
            "Virtual DOM: Efficiently updates only what changed in the UI, delivering high performance for complex interfaces",
            "Component-Based Architecture: Build encapsulated, reusable UI pieces that manage their own state and compose together",
            "Declarative Syntax (JSX): Describe what the UI should look like; React handles the how behind the scenes",
            "One-Way Data Flow: Predictable data movement from parent to child via props, making apps easier to debug",
            "Rich Ecosystem: React Router, Redux, Next.js, React Native for mobile, and thousands of community libraries",
        ],
        font_size=11, color=WHITE, bold_prefix=False
    )

    # Right side - Why React comparison box
    compare_box = slide2.shapes.add_shape(1, Inches(6.4), Inches(1.6), Inches(5.5), Inches(4.0))
    compare_box.fill.solid()
    compare_box.fill.fore_color.rgb = RGBColor(0x0F, 0x2A, 0x44)
    compare_box.line.fill.background()

    add_text_box(
        slide2, Inches(6.5), Inches(1.5), Inches(5), Inches(0.3),
        "Why React Over Other Frameworks?", font_size=13, bold=True, color=GOLD
    )

    # Comparison table as structured text
    add_text_box(
        slide2, Inches(6.6), Inches(1.95), Inches(2.5), Inches(0.25),
        "Feature", font_size=10, bold=True, color=LIGHT_BLUE
    )
    add_text_box(
        slide2, Inches(8.4), Inches(1.95), Inches(1.3), Inches(0.25),
        "React", font_size=10, bold=True, color=LIGHT_BLUE
    )
    add_text_box(
        slide2, Inches(9.6), Inches(1.95), Inches(1.2), Inches(0.25),
        "Angular", font_size=10, bold=True, color=LIGHT_BLUE
    )
    add_text_box(
        slide2, Inches(10.7), Inches(1.95), Inches(1.1), Inches(0.25),
        "Vue", font_size=10, bold=True, color=LIGHT_BLUE
    )

    # Separator line
    add_blue_accent_line(slide2, Inches(6.6), Inches(2.22), Inches(5.1))

    rows = [
        ("Type", "Library", "Framework", "Framework"),
        ("Learning Curve", "Moderate", "Steep", "Easy"),
        ("Flexibility", "Very High", "Opinionated", "Moderate"),
        ("Performance", "Virtual DOM", "Change Detect.", "Virtual DOM"),
        ("Mobile", "React Native", "Ionic", "Limited"),
        ("Community", "Largest", "Large", "Growing"),
        ("Backed By", "Meta", "Google", "Community"),
        ("Job Market", "Highest", "High", "Moderate"),
    ]

    for idx, (feat, react_val, angular_val, vue_val) in enumerate(rows):
        y_pos = Inches(2.35) + Pt(18) * idx
        row_color = WHITE if idx % 2 == 0 else RGBColor(0xCC, 0xDD, 0xEE)
        add_text_box(slide2, Inches(6.6), y_pos, Inches(2.5), Inches(0.22),
                     feat, font_size=9, bold=True, color=row_color, font_name="Roboto")
        add_text_box(slide2, Inches(8.4), y_pos, Inches(1.3), Inches(0.22),
                     react_val, font_size=9, color=RGBColor(0x7E, 0xD3, 0x21), font_name="Roboto")
        add_text_box(slide2, Inches(9.6), y_pos, Inches(1.2), Inches(0.22),
                     angular_val, font_size=9, color=row_color, font_name="Roboto")
        add_text_box(slide2, Inches(10.7), y_pos, Inches(1.1), Inches(0.22),
                     vue_val, font_size=9, color=row_color, font_name="Roboto")

    # Bottom takeaway
    add_text_box(
        slide2, Inches(6.6), Inches(4.75), Inches(5.1), Inches(0.6),
        "React's flexibility as a library (not a full framework), massive community, and seamless path to mobile via React Native make it the top choice for modern frontend development.",
        font_size=9, color=RGBColor(0xBB, 0xCC, 0xDD), font_name="Roboto"
    )

    add_slide_number(slide2, 2)

    # =========================================================================
    # SLIDE 3: Component Essentials & JSX Fundamentals
    # =========================================================================
    slide3_comp = create_content_slide(prs, 0)
    add_section_header(slide3_comp, "Component Essentials & JSX Fundamentals")

    add_bullet_list(
        slide3_comp, Inches(0.85), Inches(1.5), Inches(5.2), Inches(4.3),
        [
            "Function Components: Created reusable UI blocks using function and arrow function syntax; component names must start with a capital letter",
            "Props & Children: Passed data via props (read-only), used destructuring with default values, and composed UIs with the children prop",
            "JSX Expressions: Embedded JavaScript in JSX with {}, used ternary operators, .map() for lists with unique keys, and Fragments (<>...</>) to avoid wrapper divs",
            "Event Handlers: Attached onClick, onChange, onSubmit handlers; used arrow functions to pass parameters without immediate invocation",
            "Pure Components: Kept components pure (same input = same output); side effects belong in event handlers or useEffect, never during render",
            "Import/Export: Organized one component per file with default exports; used named exports for utility components",
        ],
        font_size=12, color=WHITE
    )

    # Right side - code snippet box
    code_box = slide3_comp.shapes.add_shape(1, Inches(6.4), Inches(1.6), Inches(5.5), Inches(3.8))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0x0F, 0x2A, 0x44)
    code_box.line.fill.background()

    add_text_box(
        slide3_comp, Inches(6.5), Inches(1.5), Inches(3), Inches(0.3),
        "Example:", font_size=11, bold=True, color=GOLD
    )

    code_text = (
        'function Greeting({ name, age }) {\n'
        '  return (\n'
        '    <div>\n'
        '      <h1>Hello, {name}!</h1>\n'
        '      <p>Age: {age}</p>\n'
        '    </div>\n'
        '  );\n'
        '}\n\n'
        'function Card({ children, title }) {\n'
        '  return (\n'
        '    <div className="card">\n'
        '      <h2>{title}</h2>\n'
        '      <div>{children}</div>\n'
        '    </div>\n'
        '  );\n'
        '}'
    )
    add_text_box(
        slide3_comp, Inches(6.6), Inches(1.9), Inches(5.1), Inches(3.3),
        code_text, font_size=10, color=RGBColor(0x7E, 0xD3, 0x21),
        font_name="Courier New"
    )
    add_slide_number(slide3_comp, 3)

    # =========================================================================
    # SLIDE 4: State, Hooks & Side Effects
    # =========================================================================
    slide4_hooks = create_content_slide(prs, 0)
    add_section_header(slide4_hooks, "State Management, Hooks & Side Effects")

    add_bullet_list(
        slide4_hooks, Inches(0.85), Inches(1.5), Inches(5.2), Inches(4.3),
        [
            "useState: Managed local component state (counters, toggles, form inputs); state updates are batched and async; used functional updates (prev => prev + 1) for previous-state-dependent logic",
            "useReducer: Handled complex state with reducer functions and dispatch actions; ideal for forms with many fields and predictable state transitions",
            "useEffect: Ran side effects after render (API calls, timers, DOM updates); used dependency arrays to control when effects fire; returned cleanup functions for subscriptions",
            "useRef: Accessed DOM elements directly (e.g., inputRef.current.focus()); stored mutable values between renders without triggering re-renders",
            "useContext: Consumed shared data (theme, auth) from Context providers; avoided prop drilling through deep component trees",
            "useMemo & useCallback: Memoized expensive calculations and function references to prevent unnecessary child re-renders",
        ],
        font_size=12, color=WHITE
    )

    # Right side - code snippet
    code_box2 = slide4_hooks.shapes.add_shape(1, Inches(6.4), Inches(1.6), Inches(5.5), Inches(3.8))
    code_box2.fill.solid()
    code_box2.fill.fore_color.rgb = RGBColor(0x0F, 0x2A, 0x44)
    code_box2.line.fill.background()

    add_text_box(
        slide4_hooks, Inches(6.5), Inches(1.5), Inches(3), Inches(0.3),
        "Example:", font_size=11, bold=True, color=GOLD
    )

    code_text2 = (
        'function Counter() {\n'
        '  const [count, setCount] =\n'
        '    useState(0);\n'
        '  const inputRef = useRef(null);\n\n'
        '  useEffect(() => {\n'
        '    document.title =\n'
        '      `Count: ${count}`;\n'
        '    return () => {\n'
        '      // cleanup\n'
        '    };\n'
        '  }, [count]);\n\n'
        '  return (\n'
        '    <>\n'
        '      <input ref={inputRef} />\n'
        '      <button onClick={() =>\n'
        '        setCount(c => c + 1)}>\n'
        '        Clicked {count} times\n'
        '      </button>\n'
        '    </>\n'
        '  );\n'
        '}'
    )
    add_text_box(
        slide4_hooks, Inches(6.6), Inches(1.9), Inches(5.1), Inches(3.3),
        code_text2, font_size=10, color=RGBColor(0x7E, 0xD3, 0x21),
        font_name="Courier New"
    )
    add_slide_number(slide4_hooks, 4)

    # =========================================================================
    # SLIDE 5: Forms, Patterns & API Integration
    # =========================================================================
    slide5_forms = create_content_slide(prs, 0)
    add_section_header(slide5_forms, "Forms, Component Patterns & API Integration")

    # Left column - three sub-sections
    add_text_box(
        slide5_forms, Inches(0.85), Inches(1.45), Inches(5), Inches(0.3),
        "Forms & Controlled Components", font_size=15, bold=True, color=GOLD
    )
    add_bullet_list(
        slide5_forms, Inches(0.85), Inches(1.85), Inches(5), Inches(1.2),
        [
            "Built controlled inputs (text, select, textarea, checkbox, radio) where React state is the single source of truth",
            "Handled form submission with e.preventDefault(); validated with HTML5 attributes and custom logic",
        ],
        font_size=11, color=WHITE, bold_prefix=False
    )

    add_text_box(
        slide5_forms, Inches(0.85), Inches(3.0), Inches(5), Inches(0.3),
        "Advanced Component Patterns", font_size=15, bold=True, color=GOLD
    )
    add_bullet_list(
        slide5_forms, Inches(0.85), Inches(3.4), Inches(5), Inches(1.2),
        [
            "Composition with children and named slot props; Compound Components sharing state via Context",
            "HOCs for cross-cutting concerns (auth, logging); Render Props for sharing logic with consumer-controlled rendering",
            "forwardRef to pass refs through components; PropTypes and TypeScript interfaces for type safety",
        ],
        font_size=11, color=WHITE, bold_prefix=False
    )

    add_text_box(
        slide5_forms, Inches(0.85), Inches(4.8), Inches(5), Inches(0.3),
        "REST API Consumption: Fetch vs Axios", font_size=15, bold=True, color=GOLD
    )
    add_bullet_list(
        slide5_forms, Inches(0.85), Inches(5.15), Inches(5), Inches(0.8),
        [
            "Fetch: built-in, two-step JSON (.json()), no cancel support; Axios: auto JSON transform, request cancellation, download progress",
            "Both consumed inside useEffect with useState for loading/error/data state management",
        ],
        font_size=11, color=WHITE, bold_prefix=False
    )

    # Right side code - LoginForm
    code_box3 = slide5_forms.shapes.add_shape(1, Inches(6.4), Inches(1.6), Inches(5.5), Inches(4.0))
    code_box3.fill.solid()
    code_box3.fill.fore_color.rgb = RGBColor(0x0F, 0x2A, 0x44)
    code_box3.line.fill.background()

    add_text_box(
        slide5_forms, Inches(6.5), Inches(1.5), Inches(3), Inches(0.3),
        "Example:", font_size=11, bold=True, color=GOLD
    )
    code_text3 = (
        'function LoginForm() {\n'
        '  const [email, setEmail] =\n'
        '    useState("");\n'
        '  const [password, setPassword] =\n'
        '    useState("");\n\n'
        '  const handleSubmit = (e) => {\n'
        '    e.preventDefault();\n'
        '    console.log(email, password);\n'
        '  };\n\n'
        '  return (\n'
        '    <form onSubmit={handleSubmit}>\n'
        '      <input type="email"\n'
        '        value={email}\n'
        '        onChange={(e) =>\n'
        '          setEmail(e.target.value)}\n'
        '      />\n'
        '      <input type="password"\n'
        '        value={password}\n'
        '        onChange={(e) =>\n'
        '          setPassword(\n'
        '            e.target.value)}\n'
        '      />\n'
        '      <button type="submit">\n'
        '        Login\n'
        '      </button>\n'
        '    </form>\n'
        '  );\n'
        '}'
    )
    add_text_box(
        slide5_forms, Inches(6.6), Inches(1.9), Inches(5.1), Inches(3.5),
        code_text3, font_size=9, color=RGBColor(0x7E, 0xD3, 0x21),
        font_name="Courier New"
    )
    add_slide_number(slide5_forms, 5)

    # =========================================================================
    # SLIDE 6: Optimization, Advanced Features & Thank You
    # =========================================================================
    slide6 = create_content_slide(prs, 0)

    # Center "THANK YOU" text
    add_text_box(
        slide6, Inches(3.5), Inches(0.5), Inches(5), Inches(0.6),
        "THANK YOU", font_size=34, bold=True, color=WHITE,
        alignment=PP_ALIGN.CENTER, font_name="Roboto"
    )
    add_blue_accent_line(slide6, Inches(4.5), Inches(1.15), Inches(3))

    add_text_box(
        slide6, Inches(1.5), Inches(1.4), Inches(9), Inches(0.4),
        "Optimization & Advanced Features Learned", font_size=18, bold=True, color=GOLD,
        alignment=PP_ALIGN.CENTER, font_name="Roboto"
    )

    # Two-column layout for advanced topics
    add_text_box(
        slide6, Inches(0.7), Inches(1.95), Inches(5), Inches(0.3),
        "Performance & Architecture", font_size=14, bold=True, color=LIGHT_BLUE
    )
    add_bullet_list(
        slide6, Inches(0.7), Inches(2.3), Inches(5.2), Inches(1.8),
        [
            "React.memo: Prevented re-renders when props are unchanged; custom comparison functions for fine control",
            "Suspense & lazy(): Code-split components into separate bundles loaded on demand with fallback UI",
            "Error Boundaries: Caught render errors in child trees with class components; used fallback UI for graceful degradation",
            "Portals: Rendered modals and tooltips outside parent DOM hierarchy while preserving React event bubbling",
        ],
        font_size=11, color=WHITE
    )

    add_text_box(
        slide6, Inches(6.3), Inches(1.95), Inches(5.5), Inches(0.3),
        "Styling & Development Tools", font_size=14, bold=True, color=LIGHT_BLUE
    )
    add_bullet_list(
        slide6, Inches(6.3), Inches(2.3), Inches(5.5), Inches(1.8),
        [
            "CSS Modules: Scoped styles without runtime cost; Styled Components for dynamic JS-powered styling",
            "Inline Styles: camelCase properties as JS objects; conditional styling with ternary expressions",
            "StrictMode: Enabled double-rendering in development to catch impure components and missing cleanup early",
            "Tailwind CSS: Utility-first approach for rapid UI development",
        ],
        font_size=11, color=WHITE
    )

    # Summary box
    summary_box = slide6.shapes.add_shape(1, Inches(0.7), Inches(4.4), Inches(11), Inches(1.3))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(0x0F, 0x2A, 0x44)
    summary_box.line.fill.background()

    add_text_box(
        slide6, Inches(0.9), Inches(4.35), Inches(10.5), Inches(0.3),
        "Skills Summary", font_size=14, bold=True, color=GOLD
    )
    add_bullet_list(
        slide6, Inches(0.9), Inches(4.7), Inches(10.5), Inches(0.9),
        [
            "Component Essentials & JSX: Function components, props, children, events, conditional rendering, lists & keys, fragments",
            "State & Hooks: useState, useReducer, useEffect, useLayoutEffect, useRef, useContext, useMemo, useCallback",
            "Forms & Patterns: Controlled inputs, form submission, composition, compound components, HOCs, render props, forwardRef",
            "API & Optimization: Fetch vs Axios, React.memo, Suspense/lazy, Error Boundaries, Portals, StrictMode, CSS Modules",
        ],
        font_size=10, color=WHITE
    )

    # Presenter info
    add_text_box(
        slide6, Inches(2.5), Inches(5.75), Inches(7), Inches(0.3),
        "NINDENKIMANA Blaise  |  Strategy and Risk Analysis Department  |  RRA",
        font_size=11, color=RGBColor(0xBB, 0xCC, 0xDD),
        alignment=PP_ALIGN.CENTER, font_name="Roboto"
    )
    add_slide_number(slide6, 6)

    prs.save(OUTPUT)
    print(f"Presentation saved to: {OUTPUT}")


if __name__ == "__main__":
    main()
